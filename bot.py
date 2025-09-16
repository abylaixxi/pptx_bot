import logging
import wikipediaapi
from pptx import Presentation
from pptx.util import Inches, Pt
from telegram import Update
from telegram.ext import ApplicationBuilder, CommandHandler, ContextTypes, MessageHandler, filters
import tempfile
import os

# --- Настройка логирования ---
logging.basicConfig(
    format="%(asctime)s - %(name)s - %(levelname)s - %(message)s", level=logging.INFO
)
logger = logging.getLogger(__name__)

# --- Настройка Wikipedia ---
wiki = wikipediaapi.Wikipedia('ru')  # Для русского, можно сменить на 'en'

# --- Создание презентации ---
def create_presentation(topic: str, slides_count: int = 8) -> str:
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # Заголовочный слайд
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = topic
    slide.placeholders[1].text = "Автоматически создано ботом"

    # Получаем текст из Wikipedia
    page = wiki.page(topic)
    if not page.exists():
        return None

    paragraphs = page.summary.split('. ')
    chunks = [paragraphs[i:i + 3] for i in range(0, len(paragraphs), 3)]

    # Создание слайдов
    for i in range(slides_count - 1):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = f"{topic} — часть {i+1}"
        text = '. '.join(chunks[i]) if i < len(chunks) else "Нет дополнительной информации."
        content.text = text

    # Сохраняем во временный файл
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        prs.save(tmp.name)
        return tmp.name

# --- Команды ---
async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text(
        "Привет! Я бот для создания презентаций.\n"
        "Напиши команду /make <тема> [кол-во слайдов]. Пример:\n"
        "/make Французская революция 8"
    )

async def make(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if len(context.args) == 0:
        await update.message.reply_text("Укажи тему. Пример: /make Французская революция 8")
        return

    # Парсим аргументы
    try:
        slides_count = int(context.args[-1])
        topic = ' '.join(context.args[:-1])
    except ValueError:
        slides_count = 8
        topic = ' '.join(context.args)

    await update.message.reply_text(f"Создаю презентацию по теме: {topic} ({slides_count} слайдов)...")

    pptx_path = create_presentation(topic, slides_count)
    if pptx_path is None:
        await update.message.reply_text("Не удалось найти информацию по этой теме.")
        return

    # Отправка файла
    with open(pptx_path, 'rb') as file:
        await update.message.reply_document(document=file, filename=f"{topic}.pptx")

    os.remove(pptx_path)

# --- Основная функция ---
def main():
    TOKEN = "ВСТАВЬ_СВОЙ_ТОКЕН_ТЕЛЕГРАМ"
    app = ApplicationBuilder().token(TOKEN).build()

    app.add_handler(CommandHandler("start", start))
    app.add_handler(CommandHandler("make", make))

    app.run_polling()

if __name__ == "__main__":
    main()
